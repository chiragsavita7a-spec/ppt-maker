"""
PPT Generator — pure python-pptx, no Node.js required.
Called by app.py: generate_slides(data) -> pptx bytes
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io, base64

# ── Themes ────────────────────────────────────────────────────────────────────
THEMES = {
    "classic": dict(
        PRIMARY=RGBColor(0xF9,0x61,0x67), SECONDARY=RGBColor(0x2F,0x3C,0x7E),
        ACCENT=RGBColor(0xF9,0xE7,0x95),  LIGHT=RGBColor(0xFD,0xDD,0xE0),
        BG=RGBColor(0xF4,0xF6,0xFB),      DARK=RGBColor(0x1A,0x1A,0x2E)),
    "ocean": dict(
        PRIMARY=RGBColor(0x00,0x77,0xB6), SECONDARY=RGBColor(0x03,0x04,0x5E),
        ACCENT=RGBColor(0x90,0xE0,0xEF),  LIGHT=RGBColor(0xCA,0xF0,0xF8),
        BG=RGBColor(0xF0,0xF7,0xFF),      DARK=RGBColor(0x03,0x04,0x5E)),
    "forest": dict(
        PRIMARY=RGBColor(0x2D,0x6A,0x4F), SECONDARY=RGBColor(0x1B,0x40,0x32),
        ACCENT=RGBColor(0xD8,0xF3,0xDC),  LIGHT=RGBColor(0xB7,0xE4,0xC7),
        BG=RGBColor(0xF0,0xFB,0xF2),      DARK=RGBColor(0x08,0x1C,0x15)),
    "royal": dict(
        PRIMARY=RGBColor(0x7B,0x2D,0x8B), SECONDARY=RGBColor(0x4A,0x00,0x72),
        ACCENT=RGBColor(0xF9,0xE7,0x95),  LIGHT=RGBColor(0xE9,0xC4,0xF0),
        BG=RGBColor(0xFB,0xF0,0xFF),      DARK=RGBColor(0x2D,0x00,0x4B)),
    "dark": dict(
        PRIMARY=RGBColor(0xF9,0x61,0x67), SECONDARY=RGBColor(0x1A,0x1A,0x2E),
        ACCENT=RGBColor(0xF9,0xE7,0x95),  LIGHT=RGBColor(0x44,0x44,0x66),
        BG=RGBColor(0x16,0x21,0x3E),      DARK=RGBColor(0x0F,0x0F,0x1A)),
    "sunrise": dict(
        PRIMARY=RGBColor(0xFF,0x6B,0x35), SECONDARY=RGBColor(0xC1,0x12,0x1F),
        ACCENT=RGBColor(0xFF,0xE0,0x66),  LIGHT=RGBColor(0xFF,0xCC,0xA1),
        BG=RGBColor(0xFF,0xF8,0xF0),      DARK=RGBColor(0x6B,0x0F,0x0F)),
    "slate": dict(
        PRIMARY=RGBColor(0x46,0x7F,0xFF), SECONDARY=RGBColor(0x1E,0x2D,0x40),
        ACCENT=RGBColor(0xA8,0xD8,0xFF),  LIGHT=RGBColor(0xD0,0xE8,0xFF),
        BG=RGBColor(0xF2,0xF5,0xFF),      DARK=RGBColor(0x0D,0x1B,0x2A)),
    "mint": dict(
        PRIMARY=RGBColor(0x00,0xB4,0x9A), SECONDARY=RGBColor(0x00,0x63,0x58),
        ACCENT=RGBColor(0xB2,0xF5,0xEA),  LIGHT=RGBColor(0xCC,0xF9,0xF4),
        BG=RGBColor(0xF0,0xFD,0xFB),      DARK=RGBColor(0x00,0x3B,0x36)),
    "crimson": dict(
        PRIMARY=RGBColor(0xC0,0x39,0x2B), SECONDARY=RGBColor(0x7B,0x24,0x1C),
        ACCENT=RGBColor(0xF9,0xE7,0x95),  LIGHT=RGBColor(0xF5,0xB7,0xB1),
        BG=RGBColor(0xFD,0xF2,0xF2),      DARK=RGBColor(0x4A,0x00,0x00)),
    "golden": dict(
        PRIMARY=RGBColor(0xF3,0x9C,0x12), SECONDARY=RGBColor(0x87,0x55,0x00),
        ACCENT=RGBColor(0xFF,0xF0,0x8A),  LIGHT=RGBColor(0xFF,0xE5,0xA0),
        BG=RGBColor(0xFF,0xFB,0xF0),      DARK=RGBColor(0x4A,0x2C,0x00)),
}
WHITE = RGBColor(0xFF,0xFF,0xFF)
OFFWH = RGBColor(0xF4,0xF6,0xFB)

# Active theme colours — swapped in by generate_slides()
_T = THEMES["classic"]
def _p(): return _T["PRIMARY"]
def _s(): return _T["SECONDARY"]
def _a(): return _T["ACCENT"]
def _l(): return _T["LIGHT"]
def _bg(): return _T["BG"]
def _dk(): return _T["DARK"]

# Legacy aliases so existing builders keep working
CORAL   = RGBColor(0xF9, 0x61, 0x67)
NAVY    = RGBColor(0x2F, 0x3C, 0x7E)
GOLD    = RGBColor(0xF9, 0xE7, 0x95)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
CORAL_L = RGBColor(0xFD, 0xDD, 0xE0)
MID     = RGBColor(0x44, 0x44, 0x66)

W = Inches(10)
H = Inches(5.625)

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font_name
    if color: run.font.color.rgb = color
    return txb

def trim(s, n):
    if not s: return ""
    return s[:n-1]+"…" if len(s)>n else s

def add_circle(slide, x, y, d, fill):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_image_from_b64(slide, b64, x, y, w, h):
    import io
    from PIL import Image
    img_bytes = base64.b64decode(b64)
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    buf = io.BytesIO(); img.save(buf, "PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

# ── Slide effects: transitions + entrance animations ─────────────────────────
def _add_slide_effects(slide, transition='push', animate=True):
    """Inject OOXML slide transition and auto-play staggered fade-in animations."""
    try:
        from lxml import etree
        P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        # 1. Slide transition
        inner = {
            'fade':  '<p:fade/>',
            'push':  '<p:push dir="l"/>',
            'wipe':  '<p:wipe dir="l"/>',
            'zoom':  '<p:zoom/>',
            'split': '<p:split orient="vert" dir="out"/>',
        }.get(transition, '<p:push dir="l"/>')
        spd = 'slow' if transition == 'zoom' else 'med'
        slide.element.append(etree.fromstring(
            f'<p:transition xmlns:p="{P}" spd="{spd}">{inner}</p:transition>'.encode()))

        if not animate:
            return

        # 2. Auto-play staggered fade-in — cap at 8 shapes to keep XML lean
        shape_ids = [sh.shape_id for sh in slide.shapes][:8]
        if not shape_ids:
            return

        parts = []
        ctn = 3
        for i, sp_id in enumerate(shape_ids):
            delay = i * 400          # 400 ms stagger
            dur   = 600              # 600 ms fade duration
            parts.append(f'''
            <p:par>
              <p:cTn id="{ctn}" fill="hold">
                <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
                <p:childTnLst>
                  <p:par>
                    <p:cTn id="{ctn+1}" presetID="10" presetClass="entr"
                           presetSubtype="0" fill="hold" grpId="{i}" nodeType="withEffect">
                      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      <p:childTnLst>
                        <p:animEffect transition="in" filter="fade">
                          <p:cBhvr>
                            <p:cTn id="{ctn+2}" dur="{dur}"/>
                            <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                          </p:cBhvr>
                        </p:animEffect>
                      </p:childTnLst>
                    </p:cTn>
                  </p:par>
                </p:childTnLst>
              </p:cTn>
            </p:par>''')
            ctn += 3

        slide.element.append(etree.fromstring(f'''<p:timing xmlns:p="{P}">
          <p:tnLst>
            <p:par>
              <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
                <p:childTnLst>
                  <p:par>
                    <p:cTn id="2" fill="hold">
                      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      <p:childTnLst>{"".join(parts)}</p:childTnLst>
                    </p:cTn>
                  </p:par>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:tnLst>
          <p:bldLst/>
        </p:timing>'''.encode()))
    except Exception:
        pass  # animations are optional — never break slide generation

# ── Slide builders ────────────────────────────────────────────────────────────
def build_title_slide(prs, data):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)

    add_rect(s, 0, 0, 10, 5.625, NAVY)
    add_rect(s, 0, 0, 0.55, 5.625, CORAL)
    add_rect(s, 0.55, 5.1, 9.45, 0.525, GOLD)

    # Decorative circle
    add_circle(s, 7.7, -1.0, 3.2, CORAL)
    add_circle(s, 8.0, -0.7, 2.7, NAVY)

    sub = f"{data.get('subject','')}{'  ·  '+data['grade'] if data.get('grade') else ''}"
    if sub.strip():
        add_textbox(s, 0.9, 0.45, 8, 0.4, trim(sub,60), 13, color=GOLD)

    add_textbox(s, 0.9, 1.0, 7.5, 2.2, trim(data.get("topic","Lesson"),70),
                42, bold=True, color=WHITE, font_name="Cambria")

    if data.get("subtitle"):
        add_textbox(s, 0.9, 3.2, 7.5, 0.7, trim(data["subtitle"],100),
                    16, italic=True, color=CORAL_L)


def build_objectives_slide(prs, data):
    if not data.get("objectives"): return
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)

    add_rect(s, 0, 0, 10, 5.625, OFFWH)
    add_rect(s, 0, 0, 0.5, 5.625, CORAL)

    add_textbox(s, 0.8, 0.25, 9, 0.75, "Learning Objectives",
                32, bold=True, color=NAVY, font_name="Cambria")
    add_rect(s, 0.8, 1.0, 8.8, 0.04, CORAL)

    colors = [CORAL, NAVY, rgb(0xE6,0x7E,0x22), rgb(0x27,0xAE,0x60), rgb(0x8E,0x44,0xAD)]
    for i, obj in enumerate(data["objectives"][:5]):
        y = 1.15 + i * 0.82
        add_circle(s, 0.8, y, 0.52, colors[i%5])
        add_textbox(s, 0.8, y, 0.52, 0.52, str(i+1), 15, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, 1.45, y+0.05, 8.0, 0.5, trim(obj,90), 14, color=DARK)


def build_content_slide(prs, slide_data, idx):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    has_img = bool(slide_data.get("image_b64"))

    add_rect(s, 0, 0, 10, 5.625, WHITE)
    add_rect(s, 0, 0, 10, 1.0, CORAL)
    add_rect(s, 9.1, 0, 0.9, 1.0, NAVY)
    add_textbox(s, 9.1, 0, 0.9, 1.0, str(idx+1), 14, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.25, 0.1, 8.7, 0.82, trim(slide_data.get("title",""),65),
                24, bold=True, color=WHITE, font_name="Cambria")

    if has_img:
        # two-column: text left, image right
        try:
            add_image_from_b64(s, slide_data["image_b64"], 6.2, 1.1, 3.4, 3.9)
        except Exception:
            add_rect(s, 6.2, 1.1, 3.4, 3.9, NAVY)
            add_rect(s, 6.2, 4.4, 3.4, 0.6, CORAL)
            kw = trim(slide_data.get("image_keyword",""), 28)
            add_textbox(s, 6.3, 4.45, 3.2, 0.5, kw, 9, italic=True, color=WHITE)

        txt_w = 5.7
        if slide_data.get("key_fact"):
            add_rect(s, 0.3, 1.15, txt_w, 1.1, NAVY)
            add_textbox(s, 0.4, 1.15, txt_w-0.1, 1.1,
                        f'"{trim(slide_data["key_fact"],90)}"',
                        13, italic=True, color=GOLD, font_name="Cambria",
                        align=PP_ALIGN.CENTER)
            start_y = 2.4
        else:
            start_y = 1.2

        for i, pt in enumerate(slide_data.get("content",[])[:5]):
            y = start_y + i * 0.72
            add_rect(s, 0.3, y+0.15, 0.1, 0.35, CORAL)
            add_textbox(s, 0.52, y, txt_w-0.22, 0.65, trim(pt,80), 13, color=DARK)

        if slide_data.get("key_term"):
            add_rect(s, 0.3, 5.08, txt_w, 0.42, GOLD)
            add_textbox(s, 0.42, 5.08, txt_w-0.12, 0.42,
                        f"Key: {trim(slide_data['key_term'],55)}", 10, bold=True, color=DARK)

    else:
        # full-width
        add_rect(s, 0.3, 1.1, 9.4, 4.25, OFFWH)
        start_y = 1.2
        if slide_data.get("key_fact"):
            add_rect(s, 0.5, 1.18, 9.0, 1.1, NAVY)
            add_textbox(s, 0.6, 1.18, 8.8, 1.1,
                        f'"{trim(slide_data["key_fact"],110)}"',
                        16, italic=True, color=GOLD, font_name="Cambria",
                        align=PP_ALIGN.CENTER)
            start_y = 2.4

        for i, pt in enumerate(slide_data.get("content",[])[:5]):
            y = start_y + i * 0.65
            add_circle(s, 0.48, y+0.08, 0.32, CORAL)
            add_textbox(s, 0.95, y, 8.7, 0.55, trim(pt,110), 14, color=DARK)

        if slide_data.get("key_term"):
            add_rect(s, 0.3, 5.1, 9.4, 0.4, GOLD)
            add_textbox(s, 0.45, 5.1, 9.1, 0.4,
                        f"Key Term: {trim(slide_data['key_term'],75)}", 10, bold=True, color=DARK)


def build_activity_slide(prs, data):
    if not data.get("activity"): return
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)

    add_rect(s, 0, 0, 10, 5.625, GOLD)
    add_rect(s, 0, 0, 10, 1.1, DARK)
    add_textbox(s, 0.3, 0.1, 9.4, 0.92, "Activity / Discussion",
                28, bold=True, color=WHITE, font_name="Cambria")
    add_rect(s, 0.5, 1.28, 9.0, 3.85, WHITE)

    # red border line at left
    add_rect(s, 0.5, 1.28, 0.06, 3.85, CORAL)

    add_textbox(s, 0.7, 1.38, 8.7, 3.65, trim(data["activity"],400),
                15, color=DARK, wrap=True)


def build_summary_slide(prs, data):
    if not data.get("summary"): return
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)

    add_rect(s, 0, 0, 10, 5.625, OFFWH)
    add_rect(s, 0, 0, 10, 1.1, NAVY)
    add_rect(s, 0, 5.2, 10, 0.425, CORAL)

    add_textbox(s, 0.3, 0.12, 9.4, 0.9, "Key Takeaways",
                28, bold=True, color=WHITE, font_name="Cambria")

    for i, pt in enumerate(data["summary"][:5]):
        y = 1.22 + i * 0.72
        add_rect(s, 0.38, y, 0.1, 0.52, CORAL)
        add_textbox(s, 0.65, y, 9.0, 0.58, trim(pt,110), 14, color=DARK)


def build_final_slide(prs, data):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)

    add_rect(s, 0, 0, 10, 5.625, NAVY)
    add_rect(s, 0, 0, 0.55, 5.625, CORAL)
    add_rect(s, 0.55, 5.1, 9.45, 0.525, GOLD)
    add_circle(s, 7.7, 3.1, 3.2, CORAL)
    add_circle(s, 8.0, 3.4, 2.7, NAVY)

    add_textbox(s, 0.9, 1.4, 7.0, 1.5, "Any Questions?",
                44, bold=True, color=WHITE, font_name="Cambria")
    add_textbox(s, 0.9, 3.0, 7.0, 0.6, trim(data.get("topic",""),70),
                16, italic=True, color=CORAL_L)
    sub = f"{data.get('subject','')}{'  ·  '+data['grade'] if data.get('grade') else ''}"
    if sub.strip():
        add_textbox(s, 0.9, 3.6, 7.0, 0.4, trim(sub,60), 13, color=GOLD)


# ── Quote / highlight break slide ────────────────────────────────────────────
def build_quote_slide(prs, slide_data):
    """Full-bleed accent slide — draws the eye between content sections."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)

    # Split background: accent left block + navy right
    add_rect(s, 0, 0, 4.2, 5.625, CORAL)
    add_rect(s, 4.2, 0, 5.8, 5.625, NAVY)

    # Large quotation mark decoration
    add_textbox(s, 0.15, -0.1, 1.5, 1.5, "“", 96, bold=True, color=WHITE,
                font_name="Georgia", align=PP_ALIGN.LEFT)

    # The key fact / quote text
    fact = trim(slide_data.get("key_fact",""), 160)
    add_textbox(s, 0.3, 1.1, 3.6, 3.0, fact, 18, italic=True, color=WHITE,
                font_name="Cambria", align=PP_ALIGN.LEFT)

    # Slide title on the right
    add_textbox(s, 4.5, 0.9, 5.1, 0.8, trim(slide_data.get("title",""),55),
                22, bold=True, color=GOLD, font_name="Cambria")
    add_rect(s, 4.5, 1.72, 5.0, 0.04, CORAL)

    # Key bullets on the right (brief)
    for i, pt in enumerate(slide_data.get("content",[])[:4]):
        y = 1.9 + i * 0.72
        add_rect(s, 4.5, y+0.2, 0.08, 0.3, GOLD)
        add_textbox(s, 4.72, y, 4.9, 0.62, trim(pt, 85), 13, color=WHITE)

    # Key term strip at bottom
    if slide_data.get("key_term"):
        add_rect(s, 4.2, 5.18, 5.8, 0.445, CORAL)
        add_textbox(s, 4.35, 5.18, 5.5, 0.445,
                    f"Key: {trim(slide_data['key_term'],60)}", 10, bold=True, color=WHITE)


# ── Main entry ────────────────────────────────────────────────────────────────
def generate_slides(data: dict) -> bytes:
    global _T, CORAL, NAVY, GOLD, DARK, CORAL_L, OFFWH
    # Apply theme
    theme_name = data.get("theme", "classic")
    _T = THEMES.get(theme_name, THEMES["classic"])
    # Remap legacy colour names to active theme so all builders use it
    CORAL   = _T["PRIMARY"]
    NAVY    = _T["SECONDARY"]
    GOLD    = _T["ACCENT"]
    DARK    = _T["DARK"]
    CORAL_L = _T["LIGHT"]
    OFFWH   = _T["BG"]

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    def _n(): return len(prs.slides)
    def _eff(t='fade', a=True):
        if prs.slides: _add_slide_effects(prs.slides[-1], t, a)

    n = _n(); build_title_slide(prs, data)
    if _n() > n: _eff('fade')

    n = _n(); build_objectives_slide(prs, data)
    if _n() > n: _eff('fade')

    slides = data.get("slides", [])
    for i, slide in enumerate(slides):
        n = _n()
        if i > 0 and i % 4 == 0 and slide.get("key_fact"):
            build_quote_slide(prs, slide)
            if _n() > n: _eff('split')
        else:
            build_content_slide(prs, slide, i)
            if _n() > n: _eff('push')

    n = _n(); build_activity_slide(prs, data)
    if _n() > n: _eff('wipe')

    n = _n(); build_summary_slide(prs, data)
    if _n() > n: _eff('fade')

    n = _n(); build_final_slide(prs, data)
    if _n() > n: _eff('fade')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
